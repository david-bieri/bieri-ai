#!/usr/bin/env python3
"""
build_kb.py  —  REAL/UAP 2004 Definitive Knowledge Base Builder
================================================================
Walks one or more folders of PPTX files, extracts every slide's title,
body text, and speaker notes, and writes a single structured Markdown
knowledge-base file suitable for uploading to a Claude Project.

Supports multiple course iterations in one KB run:
    *_S26.pptx   Spring 2026 (current)
    *_S24.pptx   Spring 2024 iteration
    *_S23.pptx   Spring 2023 iteration  (add as needed)

Usage (Windows PowerShell):
    pip install python-pptx

    # Single semester:
    python build_kb.py "C:\\...\\Slides" REALUAP2004_KB.md

    # Multiple semesters / patterns in one KB:
    python build_kb.py "C:\\...\\Slides" REALUAP2004_KB.md --pattern "*_S26.pptx,*_SU26.pptx,*_S24.pptx,*_S23.pptx"

    # Multiple root folders (e.g. slides split across years):
    python build_kb.py "C:\\...\\2026" REALUAP2004_KB.md --also "C:\\...\\2024"

    # Omit speaker notes (halves file size):
    python build_kb.py "..." REALUAP2004_KB.md --no-notes

Output slide tags:
    [NEWS]      "In the news" current-events hooks
    [LO]        Learning outcomes / key takeaways / recaps
    [DEF]       Definitions / concept introductions
    [EXAMPLE]   Worked examples, formulas, step-by-step calculations
    [ACTIVITY]  Case studies, discussion questions, in-class activities
    [CONTENT]   Standard lecture content

KB structure:
    ## Module N — Topic
    ### [S26] Lecture N.M — Title   ← vintage tag on every deck heading
    #### Slide N — [TYPE] Title
    - body bullet …
    > Notes: speaker-note text …
"""

import argparse
import re
import sys
from collections import defaultdict
from datetime import date
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("ERROR: python-pptx not found.  Run:  pip install python-pptx")


# ── Non-text asset extraction (v1.1.0) ────────────────────────────────────────
# build_kb is otherwise text-only; these helpers let the KB SEE and flag slides
# whose content is a figure, table, or chart — previously dropped silently.
def _table_to_md(tbl) -> str:
    """Render a PPTX table as a GitHub-Markdown table."""
    rows = []
    for r in tbl.rows:
        cells = [c.text.strip().replace("\n", " ").replace("|", "\\|") for c in r.cells]
        rows.append(cells)
    if not rows:
        return ""
    out = ["| " + " | ".join(rows[0]) + " |",
           "| " + " | ".join("---" for _ in rows[0]) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _chart_desc(chart) -> str:
    """One-line description of a chart: type + series + category span."""
    try:
        ctype = str(chart.chart_type).split()[0]
    except Exception:
        ctype = "chart"
    series = []
    try:
        series = [s.name for s in chart.series if s.name]
    except Exception:
        pass
    cats = []
    try:
        cats = [str(c) for c in chart.plots[0].categories if c is not None]
    except Exception:
        pass
    parts = [ctype]
    if series:
        parts.append("series: " + ", ".join(series[:4]) + ("…" if len(series) > 4 else ""))
    if cats:
        span = f"{cats[0]}–{cats[-1]}" if len(cats) > 1 else cats[0]
        parts.append(f"categories: {span}")
    return " · ".join(parts)


def extract_assets(slide) -> dict:
    """Detect pictures, tables, and charts that text extraction misses."""
    pics, tables, charts = 0, [], []
    for sh in slide.shapes:
        try:
            if getattr(sh, "has_table", False) and sh.has_table:
                tables.append(_table_to_md(sh.table))
            elif getattr(sh, "has_chart", False) and sh.has_chart:
                charts.append(_chart_desc(sh.chart))
            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
        except Exception:
            pass
    return {"pictures": pics, "tables": tables, "charts": charts}


# ── Vintage detection ─────────────────────────────────────────────────────────
# Maps filename suffix fragment → human label shown in KB headings.
# Add new entries here as older iterations surface.
VINTAGE_MAP = {
    "S26":  "Spring 2026",  "SU26": "Summer 2026",
    "S25":  "Spring 2025",  "SU25": "Summer 2025",
    "S24":  "Spring 2024",  "SU24": "Summer 2024",
    "24":   "2024",
    "S23":  "Spring 2023",  "SU23": "Summer 2023",  "23": "2023",
    "S22":  "Spring 2022",  "SU22": "Summer 2022",  "22": "2022",
    "S21":  "Spring 2021",  "21":   "2021",
    "S20":  "Spring 2020",  "20":   "2020",
}

def detect_vintage(stem: str) -> str:
    """Return a short vintage label from the filename stem."""
    # Match trailing _TOKEN (case-insensitive)
    m = re.search(r'_([A-Za-z0-9]+)$', stem)
    if m:
        key = m.group(1).upper()
        return VINTAGE_MAP.get(key, key)
    return "unknown"


# ── Slide-type heuristics ─────────────────────────────────────────────────────
_NEWS_RE     = re.compile(r'in the.{0,20}news|micro.?news|macro.?news|policy.?news', re.I)
_OUTCOMES_RE = re.compile(r'learning outcome|key learning|key take.?away|take.?away|recap|objectives', re.I)
_DEF_RE      = re.compile(r'\bwhat is\b|definition|key concept|glossary', re.I)
_EXAMPLE_RE  = re.compile(
    r'worked|example\s*\d|step \d+|formula|calculate|lq\s*=|ebm\s*=|'
    r'δ\s*(in|total|household|resident)|change in (total|basic|household)', re.I)
_ACTIVITY_RE = re.compile(r'case study|activity|discussion|in.?class|exercise', re.I)

def slide_type(title: str, body: str) -> str:
    combined = title + " " + body
    if _NEWS_RE.search(title):           return "NEWS"
    if _OUTCOMES_RE.search(title):       return "LO"
    if _DEF_RE.search(title):            return "DEF"
    if _EXAMPLE_RE.search(combined):     return "EXAMPLE"
    if _ACTIVITY_RE.search(combined):    return "ACTIVITY"
    return "CONTENT"


# ── Text extraction ───────────────────────────────────────────────────────────
def extract_slide(slide):
    """Return (title, body_lines, notes_text)."""
    texts = [
        para.text.strip()
        for sh in slide.shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs
        if para.text.strip()
    ]
    title      = texts[0] if texts else ""
    body_lines = [t for t in texts[1:] if t != title]

    notes = ""
    try:
        if slide.has_notes_slide:
            ntf = slide.notes_slide.notes_text_frame
            if ntf is not None:
                raw = ntf.text.strip()
                raw = re.sub(r'^(Instructor Note|Instuctor Note|Notes?):\s*', '', raw, flags=re.I)
                notes = re.sub(r'\s+', ' ', raw).strip()
    except Exception:
        pass

    return title, body_lines, notes


# ── Filename → (module_num, lecture_label, deck_label) ───────────────────────
def parse_stem(stem: str) -> tuple[str, str, str]:
    """
    Handles both naming conventions found in the corpus:
        2.1-Value_S26            → mod "2", lec "2.1", label "Value"
        10.2-REAL-UrbSpatStr_S24 → mod "10", lec "10.2", label "Urb Spat Str"
        2_2-5Attributes_S26      → mod "2", lec "2.2", label "5 Attributes"
        10_2_CommercialLeases_24 → mod "10", lec "10.2", label "Commercial Leases"
        review_S26               → mod "0", lec "?", label "Review"
    """
    clean = re.sub(r'_[A-Za-z0-9]+$', '', stem)   # strip vintage suffix
    # Pattern A: N.M-Label  (dot separator — actual corpus format)
    m = re.match(r'^(\d+)\.(\d+)[\-_](.+)$', clean)
    if m:
        mod, lec_n, label = m.group(1), m.group(2), m.group(3)
    else:
        # Pattern B: N_M-Label or N_M_Label  (underscore separator)
        m = re.match(r'^(\d+)[_\-](\d+)[_\-](.+)$', clean)
        if m:
            mod, lec_n, label = m.group(1), m.group(2), m.group(3)
        else:
            label = re.sub(r'([A-Z])', r' \1', clean).strip()
            return "0", "?", re.sub(r'\s+', ' ', label).strip()
    # Clean up label: strip REAL- prefix, expand CamelCase
    label = re.sub(r'^REAL[\-_]', '', label, flags=re.I)
    label = re.sub(r'([a-z])([A-Z])', r'\1 \2', label)  # camelCase
    label = re.sub(r'([0-9])([A-Za-z])', r'\1 \2', label)  # digit-letter
    label = label.replace('-', ' ').replace('_', ' ')
    label = re.sub(r'\s+', ' ', label).strip()
    return mod, f"{mod}.{lec_n}", label


# ── Vintage-aware module name tables (keyed by VINTAGE_MAP values) ──────────
# S26 module structure (13 instructional modules + review)
_MOD_S26 = {
    "0":  "General / Unclassified",
    "1":  "Introduction to Real Estate",
    "2":  "Real Estate Value Drivers",
    "3":  "Real Estate Economics",
    "4":  "Real Estate Law",
    "5":  "Time Value of Money",
    "6":  "Introduction to Appraisal",
    "7":  "Buying Residential Real Estate",
    "8":  "Mortgage Fundamentals",
    "9":  "Buying Commercial Real Estate",
    "10": "REITs",
    "11": "Sustainability Strategy & Initiatives",
    "12": "Real Estate Development",
    "13": "Public Policy Impact on Minority Communities",
    "14": "Housing Policy & Public Finance",
    "15": "Course Overview & Review",
}

# S24 module structure (8 thematic sections mapped to lecture-number blocks)
# Section boundaries: I=1-2, II=2-5, III=6-7, IV=8-10, V=10-11, VI=12-13, VII=14-15
_MOD_S24 = {
    "0":  "General / Unclassified",
    "1":  "Introduction & Review",
    "2":  "Real Estate Markets & Institutions — I",
    "3":  "Real Estate Markets & Institutions — II (Housing Cycles, Land Leverage)",
    "4":  "Real Estate Law & Property Rights",
    "5":  "Community Planning, Entitlements & Financing",
    "6":  "Real Estate Service Industry & Leases",
    "7":  "Brokerage, Appraisal & Commercial RE — I",
    "8":  "Commercial & Industrial Real Estate — II",
    "9":  "Property & Asset Management",
    "10": "Urban Spatial Structure",
    "11": "Urban & Regional Economics",
    "12": "Real Estate Finance — Risk, Return & TVM",
    "13": "Mortgage Mechanics & Cash Flow Analysis",
    "14": "Feasibility Analysis, DCF & Housing Policy",
    "15": "Review",
}

# Fallback for any vintage not explicitly mapped
_MOD_DEFAULT = _MOD_S26

# Vintages that use the S24 structure
_S24_VINTAGES = {"Spring 2024", "Spring 2023", "2024", "2023"}

def get_module_name(module_num: str, vintage: str) -> str:
    table = _MOD_S24 if vintage in _S24_VINTAGES else _MOD_S26
    return table.get(module_num, f"Module {module_num}")


# ── Collect decks from roots + patterns ───────────────────────────────────────
def collect_decks(roots: list[Path], patterns: list[str]) -> list[Path]:
    found = []
    for root in roots:
        for pat in patterns:
            found.extend(
                d for d in root.rglob(pat)
                if not d.name.startswith("~$")
            )
    # Deduplicate, sort by stem for reproducible output
    seen = set()
    unique = []
    for d in sorted(found, key=lambda p: p.stem.lower()):
        if d.resolve() not in seen:
            seen.add(d.resolve())
            unique.append(d)
    return unique


# ── KB builder ────────────────────────────────────────────────────────────────
def build_kb(decks: list[Path], include_notes: bool) -> str:
    lines = []

    # ── Header ──
    lines += [
        "# REAL/UAP 2004 — Principles of Real Estate: Course Knowledge Base",
        "",
        f"*Generated {date.today().isoformat()} · {len(decks)} decks · "
        f"Prof. David Bieri · Virginia Tech SPIA*",
        "",
        "---",
        "",
        "## How to Use This KB",
        "",
        "Each slide is tagged by type: **[NEWS]** current-events hooks, "
        "**[LO]** learning outcomes/recaps, **[DEF]** definitions, "
        "**[EXAMPLE]** worked examples and formulas, "
        "**[ACTIVITY]** case studies and discussion questions, "
        "**[CONTENT]** standard lecture content.",
        "",
        "Each deck heading carries a **vintage tag** — e.g. `[Spring 2026]` or `[2024]` — "
        "so you can distinguish material from different course iterations. "
        "Where the same concept appears in multiple vintages, both versions are present.",
        "",
        "---",
    ]

    # ── Vintage summary ──
    vintages: dict[str, int] = defaultdict(int)
    for d in decks:
        vintages[detect_vintage(d.stem)] += 1
    lines += ["", "## Included Course Iterations", ""]
    for v, n in sorted(vintages.items()):
        lines.append(f"- **{v}** — {n} deck(s)")
    lines += ["", "---"]

    # ── Curriculum structure note ──
    has_s24 = any(v in _S24_VINTAGES for v in vintages.keys())
    has_s26 = any(v not in _S24_VINTAGES for v in vintages.keys())
    if has_s24 and has_s26:
        lines += [
            "",
            "## ⚠️  Module Numbering Note",
            "",
            "The two course iterations use the **same lecture numbers (1.x–15.x) but different "
            "module content**. Module 3 in Spring 2024 covers housing cycles and land leverage; "
            "Module 3 in Spring 2026 covers real estate economics. Each deck heading carries a "
            "vintage tag so content from each iteration is always identifiable. "
            "Module headings in this KB reflect the correct topic for each vintage.",
            "",
            "| Module | Spring 2026 | Spring 2024 |",
            "|--------|-------------|-------------|",
        ]
        for n in [str(i) for i in range(1, 16)]:
            s26 = _MOD_S26.get(n, "—")
            s24 = _MOD_S24.get(n, "—")
            marker = " ⚠️" if s26 != s24 else ""
            lines.append(f"| {n} | {s26} | {s24}{marker} |")
        lines += ["", "---"]
    else:
        lines += ["", "---"]

    # ── Group by module ──
    by_module: dict[str, list[tuple[str, str, str, Path]]] = defaultdict(list)
    for deck in decks:
        mod, lec, label = parse_stem(deck.stem)
        vintage = detect_vintage(deck.stem)
        by_module[mod].append((lec, label, vintage, deck))

    total_slides = 0
    errors = []

    for mod in sorted(by_module.keys(), key=lambda x: int(x) if x.isdigit() else 999):
        vintages_in_mod = sorted({v for _, _, v, _ in by_module[mod]})
        names = list(dict.fromkeys(get_module_name(mod, v) for v in vintages_in_mod))
        if len(names) == 1:
            mod_header = names[0]
        else:
            # Two vintages with different names for same module number — show both
            mod_header = "  |  ".join(
                f"[{v.split()[0][:2]}{v.split()[-1][-2:]}] {get_module_name(mod, v)}"
                for v in vintages_in_mod
            )
        lines += ["", f"## Module {mod} — {mod_header}", ""]

        for lec, label, vintage, deck_path in sorted(by_module[mod], key=lambda x: (x[0], x[2])):
            try:
                prs = Presentation(str(deck_path))
            except Exception as e:
                msg = f"⚠️  Could not open `{deck_path.name}`: {e}"
                lines.append(f"\n> {msg}\n")
                errors.append(msg)
                continue

            n_slides = len(prs.slides)
            total_slides += n_slides
            lines += [
                f"### [{vintage}] Lecture {lec} — {label}",
                f"*`{deck_path.name}` · {n_slides} slides*",
                "",
            ]

            for slide_num, slide in enumerate(prs.slides, 1):
                title, body_lines, notes = extract_slide(slide)
                assets = extract_assets(slide)
                has_assets = assets["pictures"] or assets["tables"] or assets["charts"]
                if not title and not body_lines and not has_assets:
                    continue   # truly blank slide

                # Slide type: text heuristics first; else classify by dominant asset.
                if title or body_lines:
                    stype = slide_type(title, " ".join(body_lines))
                elif assets["tables"]:
                    stype, title = "TABLE", title or "(table)"
                elif assets["charts"]:
                    stype, title = "CHART", title or "(chart)"
                else:
                    stype, title = "FIGURE", title or "(image-only slide)"

                lines.append(f"#### Slide {slide_num} — [{stype}] {title}")

                for bl in body_lines:
                    if bl.strip() and bl.strip() != title.strip():
                        lines.append(f"- {bl}")
                if body_lines:
                    lines.append("")

                # Surface non-text assets so KB-dependent skills can see them.
                for tbl_md in assets["tables"]:
                    if tbl_md:
                        lines += ["", "**[TABLE]**", "", tbl_md, ""]
                for ch in assets["charts"]:
                    lines.append(f"> **[CHART]** {ch}")
                if assets["pictures"]:
                    lines.append(f"> **[FIGURE]** {assets['pictures']} image(s) — "
                                 f"visual content not captured as text")
                if has_assets:
                    lines.append("")

                if include_notes and notes:
                    preview = notes if len(notes) <= 500 else notes[:497] + "…"
                    lines.append(f"> **Notes:** {preview}")
                    lines.append("")

    # ── Footer ──
    lines += [
        "---",
        "",
        f"*KB totals: {total_slides} slides across {len(decks)} decks.*",
    ]
    if errors:
        lines += ["", "### ⚠️  Extraction errors", ""]
        for e in errors:
            lines.append(f"- {e}")

    return "\n".join(lines)


# ── CLI ───────────────────────────────────────────────────────────────────────
def main():
    p = argparse.ArgumentParser(
        description="Build a Markdown KB from REAL/UAP 2004 PPTX decks (multi-vintage).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument("root",   help="Primary folder to search for PPTX files")
    p.add_argument("output", help="Output .md filename (e.g. REALUAP2004_KB.md)")
    p.add_argument(
        "--pattern", default="*_S26.pptx,*_SU26.pptx,*_S24.pptx,*_S23.pptx",
        help="Comma-separated glob patterns (default: '*_S26.pptx,*_SU26.pptx,*_S24.pptx,*_S23.pptx')"
    )
    p.add_argument(
        "--also", action="append", default=[], metavar="FOLDER",
        help="Additional root folder(s) to search (repeatable)"
    )
    p.add_argument(
        "--no-notes", dest="notes", action="store_false",
        help="Omit speaker notes (halves output size)"
    )
    p.set_defaults(notes=True)
    args = p.parse_args()

    roots = [Path(args.root).expanduser().resolve()]
    for extra in args.also:
        roots.append(Path(extra).expanduser().resolve())

    for r in roots:
        if not r.is_dir():
            sys.exit(f"ERROR: not a directory: {r}")

    patterns = [pat.strip() for pat in args.pattern.split(",") if pat.strip()]

    print(f"Searching {len(roots)} folder(s) for patterns: {patterns}")
    decks = collect_decks(roots, patterns)

    if not decks:
        sys.exit(f"No matching files found. Check --pattern and folder paths.")

    print(f"\nFound {len(decks)} deck(s):")
    for d in decks:
        vintage = detect_vintage(d.stem)
        print(f"  [{vintage:>12s}]  {d.name}")

    print("\nExtracting …")
    kb_text = build_kb(decks, include_notes=args.notes)

    out = Path(args.output)
    out.write_text(kb_text, encoding="utf-8")
    size_kb = out.stat().st_size // 1024

    print(f"\n✓ KB written → {out.resolve()}")
    print(f"  {len(decks)} decks  ·  size: ~{size_kb} KB")
    if size_kb > 4000:
        print("  ⚠  Approaching 5 MB Project limit.")
        print("     Try --no-notes, or split by vintage (run twice with different --pattern).")
    else:
        print("  ✓ Within Claude Project 5 MB file limit.")


if __name__ == "__main__":
    main()
